from pathlib import Path

import fitz  # PyMuPDF
import pandas as pd
from pptx import Presentation

from core.exceptions import ParseError
from core.ingestion.ocr import extract_text_from_image


def parse_pdf(file_path: Path) -> str:
    text_content = []
    try:
        doc = fitz.open(file_path)
        for page in doc:
            page_text = page.get_text("text").strip()
            if not page_text:
                pix = page.get_pixmap()
                image_bytes = pix.tobytes("png")
                page_text = extract_text_from_image(image_bytes)

            if page_text:
                text_content.append(page_text)
    except ParseError:
        raise
    except Exception as e:
        raise ParseError(f"Error parsing PDF {file_path}: {e}") from e

    return "\n\n".join(text_content)


def parse_pptx(file_path: Path) -> str:
    text_content = []
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            text = "\n".join(slide_text).strip()
            if text:
                text_content.append(text)
    except Exception as e:
        raise ParseError(f"Error parsing PPTX {file_path}: {e}") from e

    return "\n\n".join(text_content)


def parse_xlsx(file_path: Path) -> str:
    text_content = []
    try:
        df_dict = pd.read_excel(file_path, sheet_name=None, engine="openpyxl")
        for sheet_name, df in df_dict.items():
            text_content.append(f"Sheet: {sheet_name}")
            text_content.append(df.to_csv(index=False))
    except Exception as e:
        raise ParseError(f"Error parsing XLSX {file_path}: {e}") from e

    return "\n\n".join(text_content)


def parse_document(file_path: Path) -> str:
    ext = file_path.suffix.lower()
    if ext == ".pdf":
        return parse_pdf(file_path)
    if ext == ".pptx":
        return parse_pptx(file_path)
    if ext in [".xlsx", ".xls"]:
        return parse_xlsx(file_path)
    return ""
