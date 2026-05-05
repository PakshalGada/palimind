from dataclasses import dataclass
from pathlib import Path
from typing import Iterator

MAX_FILE_BYTES = 1_000_000

DEFAULT_EXTENSIONS = {
    ".txt",
    ".md",
    ".py",
    ".js",
    ".ts",
    ".go",
    ".rs",
    ".c",
    ".cpp",
    ".h",
    ".json",
    ".yaml",
    ".yml",
    ".toml",
    ".csv",
    ".html",
    ".rst",
    ".pdf",
    ".docx",
    ".pptx",
}


@dataclass
class FileEntry:
    absolute: Path
    relative: Path
    size: int


def _has_hidden_components(path: Path, root: Path) -> bool:
    try:
        rel = path.relative_to(root)
    except ValueError:
        return False
    return any(part.startswith(".") for part in rel.parts[:-1])


def _extract_text(filepath: Path) -> str | None:
    suffix = filepath.suffix.lower()
    try:
        if suffix == ".pdf":
            import fitz  # pymupdf

            doc = fitz.open(str(filepath))
            return "\n".join(page.get_text() for page in doc)
        elif suffix == ".docx":
            import docx

            doc = docx.Document(str(filepath))
            return "\n".join(p.text for p in doc.paragraphs)
        elif suffix == ".pptx":
            from pptx import Presentation

            prs = Presentation(str(filepath))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            return "\n".join(texts)
        else:
            return filepath.read_text(encoding="utf-8")
    except Exception:
        return None


def walk(root: Path, allowed_extensions: set[str] | None = None) -> Iterator[FileEntry]:
    """Yield FileEntry for every valid file under root."""
    extensions = (
        allowed_extensions if allowed_extensions is not None else DEFAULT_EXTENSIONS
    )

    for filepath in root.rglob("*"):
        if not filepath.is_file():
            continue

        if _has_hidden_components(filepath, root):
            continue

        if filepath.suffix.lower() not in extensions:
            continue

        try:
            size = filepath.stat().st_size
        except OSError:
            continue
        if size > MAX_FILE_BYTES:
            continue

        if _extract_text(filepath) is None:
            continue

        yield FileEntry(
            absolute=filepath.resolve(),
            relative=filepath.relative_to(root),
            size=size,
        )
